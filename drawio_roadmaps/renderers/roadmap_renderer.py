from drawio_roadmaps.drawio.drawio_shapes import Rectangle
from drawio_roadmaps.drawio.drawio_utils import id_generator_2
from drawio_roadmaps.enums.event_type import EventType
from drawio_roadmaps.enums.swimlane_type import SwimlaneType
from drawio_roadmaps.drawio import drawio_utils

from drawio_roadmaps.config import DRAWIO_EXECUTABLE_PATH
from drawio_roadmaps.config import RoadmapConfig as config

from datetime import datetime, date
from lxml import etree
import os


class RoadmapRenderer:
    def render(self, roadmap):
        raise NotImplementedError


class StringRoadmapRenderer():
    def render(self, roadmap):
        return str(roadmap)


class AsciiRoadmapRenderer():
    def __init__(self):
        # The quarters don't line up properly so extend if showing quarter
        # also likely to have more dense data so this is likely beneficial anyway
        if config.Global.show_quarters:
            self.segment_width = config.Ascii.segment_width - 1
        else:
            self.segment_width = config.Ascii.segment_width

    def render(self, roadmap):
        result = ""

        # Draw the title and legend
        pad = "#" * (self.segment_width + 2) + '\n'
        title = f"Roadmap: {roadmap.name}\n"
        title_and_legend = pad + title + pad + EventType.get_legend_for_event_types() + '\n' \
                           + SwimlaneType.get_legend_for_swimlane_types() + '\n'
        result += title_and_legend

        # Draw the header
        horizontal_segment = '+' + '-' * self.segment_width
        roadmap_delim = horizontal_segment * (roadmap.years + 1) + '+\n'
        result += roadmap_delim

        header_segment = '| ' + ' ' * (self.segment_width - 1)
        roamap_pad = header_segment * (roadmap.years + 1) + '|\n'
        result += roamap_pad

        # put columne title in first segment and pad till segment width
        roadmap_header_text = '| ' + roadmap.swimlane_column_title
        roadmap_header_text += ' ' * (self.segment_width - len(roadmap_header_text)) + ' |'

        # in second and following segments, center the text, which will be the years
        for year in range(roadmap.start_year, roadmap.end_year):
            year_str = f"{year}"
            year_str = year_str.center(self.segment_width)
            roadmap_header_text += year_str + '|'
        result += roadmap_header_text + '\n'

        if config.Global.show_quarters:
            result += roamap_pad
            result += '| ' + ' ' * (self.segment_width - 1) + '|'
            # Draw the quarter markers
            for year in range(roadmap.years):
                for quarter in range(1, 5):
                    quarter_marker = f"Q{quarter}"
                    result += '   ' + quarter_marker + '   |'

            result += '\n'
        else:
            result += roamap_pad
        result += roadmap_delim

        for swimlane in roadmap.swimlanes:

            result += swimlane.render()
            result += roadmap_delim

        return result


class DrawIORoadmapRenderer:

    def render(self, roadmap):

        year_lenght_px = config.DrawIO.year_length_px
        swimlane_height_px = config.DrawIO.swimlane_height_px

        mxGraphModel = drawio_utils.get_diagram_root()
        root = mxGraphModel.find("root")
        self.append_layers(root)

        # Todo get years based on dates in roadmap - this is not implemented in Ascii
        years = [str(x) for x in range(roadmap.start_year, roadmap.start_year + roadmap.years)]

        # write out the header
        self.create_header(roadmap, root, swimlane_height_px, year_lenght_px, years)


        xy_cursor = (0, swimlane_height_px) # allow for header

        for index, swimlane in enumerate(roadmap.swimlanes):

            # the + allows for the header
            # the idea with xy_cursor is to keep track of where we are in the diagram,
            # but it needs to account for variable additional lifelines in each swimlane
            # so we need to work out the total height of each swimlane ahead of time
            # and then use that to calculate the y position of the next swimlane


            lane = Rectangle('',
                             xy_cursor[0] + year_lenght_px,
                             xy_cursor[1],
                             year_lenght_px * (len(years)),
                             swimlane.height(),
                             style={
                                    'strokeColor': '#000000;',
                                    'fontStyle': '1',
                                    })
            lane.render(root)

            typographic_line_gap = 20  # todo remove magic

            style = {'fontStyle': '1'}
            # change the style of the swimlane titles (for all) if we have any lifelines in the roadmap
            if any(swimlane.lifelines for swimlane in roadmap.swimlanes):
                style.update({'verticalAlign': 'top',
                              'spacingTop': '36',
                              'spacingLeft': typographic_line_gap,
                              'align': 'left'})

            swimlane_label = Rectangle(swimlane.name,
                                       xy_cursor[0],
                                       xy_cursor[1],
                                       year_lenght_px,
                                       swimlane.height(),
                                       style=style)
            swimlane_label.render(root)



            xy_timeline_begin = (xy_cursor[0] + year_lenght_px + typographic_line_gap
                                 , xy_cursor[1] + int(swimlane_height_px / 2))

            xy_timeline_end = (
            xy_cursor[0] + year_lenght_px + year_lenght_px * (len(years)) - typographic_line_gap
            , xy_cursor[1] + int(swimlane_height_px / 2))

            # because we now have optional from_date and to_date on the swimlane class we need to work out the start
            # and end of the timeline for each swimlane in x and y position and then pass that to the lifeline renderer
            # we also need to check if we have a merge_to lifeline and if so, render that as an angled line

            if swimlane.date_from:
                xy_timeline_begin = (xy_timeline_begin[0] + year_lenght_px * (swimlane.date_from.year - roadmap.start_year), xy_timeline_begin[1])

            if swimlane.date_to:
                xy_timeline_end = (xy_timeline_end[0] - year_lenght_px * (roadmap.end_year - swimlane.date_to.year), xy_timeline_end[1])


            swimlane.tubemap_line(root=root,
                                  layer="Default",
                                  begin_x=xy_timeline_begin[0],
                                  begin_y=xy_timeline_begin[1],
                                  end_x=xy_timeline_end[0],
                                  end_y=xy_timeline_end[1],
                                  width=2,
                                  height=2,
                                  style={
                                      'strokeColor': swimlane.type.metadata_drawio.strokeColor,
                                      'strokeWidth': '5',
                                      'endArrow': 'doubleBlock',
                                      'endSize': '1'
                                  },
                                  value='')

            for event in swimlane.events:
                x = xy_cursor[0] + \
                    year_lenght_px + \
                    year_lenght_px * (event.date.year - roadmap.start_year) + \
                    year_lenght_px / 12 * (event.date.month - 1)

                # ToDo: hardcoded half circle height @ 9px (total size 18px)
                # should go to config, although unlikely to change
                # so 9 in the line below for half tube station height

                y = xy_cursor[1] + int(swimlane_height_px / 2) - 9
                event.tubemap_station(root=root,
                                      layer="Default",
                                      x=x,
                                      y=y,
                                      style={
                                          'fillColor': event.event_type.render_meta.fillColor,
                                      })

            xy_timeline_begin = (xy_timeline_begin[0] - typographic_line_gap, xy_timeline_begin[1] + swimlane_height_px // 2 + 10) # todo remove magic
            xy_timeline_end = (xy_timeline_end[0] - typographic_line_gap, xy_timeline_end[1] + swimlane_height_px // 2 + 10)

            delayed_render_lifelines = []

            for ix_lf, lifeline in enumerate(swimlane.lifelines):

                start_gap = typographic_line_gap if (lifeline.date_from is None or
                                                     lifeline.date_from == date(roadmap.start_year,1,1)) \
                                                 else 0
                lifeline.date_from = lifeline.date_from or date(roadmap.start_year, 1, 1)

                # Todo: change styling of the the terminating arrow as well
                end_gap = typographic_line_gap if (lifeline.date_to is None or
                                                   lifeline.date_to == date(roadmap.end_year-1, 12, 31)) \
                                               else 0
                if end_gap:
                    end_arrow_style = {'endArrow': 'doubleBlock', 'endSize': '1'}

                lifeline.date_to = lifeline.date_to or date(roadmap.start_year + roadmap.years, 1, 1)


                start_position_ratio = (lifeline.date_from - date(roadmap.start_year, 1, 1)).days / (365.25 * roadmap.years)
                end_position_ratio = (lifeline.date_to - date(roadmap.start_year, 1, 1)).days / (365.25 * roadmap.years)

                lifeline_begin_x = xy_timeline_begin[0] + start_position_ratio * year_lenght_px * roadmap.years + start_gap
                lifeline_end_x = xy_timeline_begin[0] + end_position_ratio * year_lenght_px * roadmap.years - end_gap

                # ToDo: why are we passing so much when it's on the object?
                # it should be just lifeline.render() with positional and layer data?

                if not lifeline.merge_to:

                    lifeline.tubemap_lifeline(root=root,
                                              layer="Default",
                                              begin_x=lifeline_begin_x,
                                              begin_y=xy_timeline_begin[1] + ix_lf * (swimlane_height_px // 4),  # todo remove magic
                                              end_x=lifeline_end_x,
                                              end_y=xy_timeline_end[1] + ix_lf * (swimlane_height_px // 4),  # todo remove magic
                                              width=2,
                                              height=2,
                                              style={
                                                'strokeColor': lifeline.type.metadata_drawio.strokeColor,
                                                'strokeWidth': '5',
                                                'endArrow': 'oval',
                                                } | (end_arrow_style if end_gap else {}),
                                              value=lifeline.name)

                else:

                    merge_to_y = swimlane.get_lifeline_y_coordinate_index(lifeline.merge_to)
                    if merge_to_y is not None:
                        lf = lifeline.tubemap_lifeline_angled(root=root,
                                                  layer="Default",
                                                  begin_x=lifeline_begin_x,
                                                  begin_y=xy_timeline_begin[1] + ix_lf * (swimlane_height_px // 4),  # todo remove magic
                                                  end_x=lifeline_end_x,
                                                  end_y= xy_timeline_end[1] + merge_to_y * (swimlane_height_px // 4),  # todo remove magic
                                                  width=2,
                                                  height=2,
                                                  style={
                                                      'strokeColor': lifeline.type.metadata_drawio.strokeColor,
                                                      'strokeWidth': '5',
                                                      'endArrow': 'oval',
                                                  } | (end_arrow_style if end_gap else {}),
                                                  value=lifeline.name)
                        delayed_render_lifelines.append(lf)


                lifeline.tubemap_lifeline_label(root=root,
                                                x=0 + typographic_line_gap,
                                                y=xy_timeline_begin[1] + ix_lf * (swimlane_height_px // 4) - 10,
                                                width=year_lenght_px - typographic_line_gap,
                                                height=20,
                                                value=lifeline.name[:36], # TODO This is truncation magic should be refactored
                                                style = {
                                                    'fontSize': '12',
                                                    'align': 'left',
                                                    'fontColor': lifeline.type.metadata_drawio.strokeColor,
                                                    'fontStyle': '1',
                                                })

                for lf in delayed_render_lifelines:
                    lf.render(root)

            print(f"Swimlane {swimlane.name} height: {swimlane.height()}")
            actual_height = swimlane.height()
            xy_cursor = (xy_cursor[0], xy_cursor[1] + actual_height)

        # "Pretty Print" to console is not really required but we like to pretty print the XML just for comparison
        # and visual confirmation of what's being produced
        # Encoding required for use in Confluence/Web
        # Open desktop drawio for convenience

        drawio_utils.pretty_print_to_console(mxGraphModel)
        drawio_utils.encode_and_save_to_file(mxGraphModel)
        os.system(f'"{DRAWIO_EXECUTABLE_PATH}" output.drawio')

        return "DrawIO roadmap generated successfully!"

    def create_header(self, roadmap, root, swimlane_height, year_lenght_px, years):
        swimlane_header = Rectangle(roadmap.swimlane_column_title, 0, 0, year_lenght_px, swimlane_height,
                                    style={
                                        'fillColor': '#ffffff',
                                        'strokeColor': '#000000;',
                                        'fontStyle': '1',
                                    })
        swimlane_header.render(root)

        for index, year in enumerate(years):
            xy_cursor = (year_lenght_px * (index + 1), 0)
            year_label = Rectangle(year,
                                   x=xy_cursor[0],
                                   y=xy_cursor[1],
                                   width=year_lenght_px,
                                   height=swimlane_height * 0.75 if config.Global.show_quarters else swimlane_height,
                                   style={
                                       'fillColor': '#ffffff',
                                       'strokeColor': '#000000;',
                                       'fontStyle': '1',
                                   })
            year_label.render(root)



        if config.Global.show_quarters:
            style = {
                     'fillColor': '#ffffff',
                     'strokeColor': '#000000;',
                     'fontSize': '12',
                    }
            # Copilot on fire
            for index, year in enumerate(years):
                xy_cursor = (year_lenght_px * (index + 1), swimlane_height * 0.75)
                quarter_label = Rectangle(f"Q1",
                                          x=xy_cursor[0],
                                          y=xy_cursor[1],
                                          width=year_lenght_px / 4,
                                          height=swimlane_height * 0.25,
                                          style=style)
                quarter_label.render(root)
                xy_cursor = (year_lenght_px * (index + 1) + year_lenght_px / 4, swimlane_height * 0.75)
                quarter_label = Rectangle(f"Q2",
                                          x=xy_cursor[0],
                                          y=xy_cursor[1],
                                          width=year_lenght_px / 4,
                                          height=swimlane_height * 0.25,
                                          style=style)
                quarter_label.render(root)
                xy_cursor = (year_lenght_px * (index + 1) + year_lenght_px / 2, swimlane_height * 0.75)
                quarter_label = Rectangle(f"Q3",
                                          x=xy_cursor[0],
                                          y=xy_cursor[1],
                                          width=year_lenght_px / 4,
                                          height=swimlane_height * 0.25,
                                          style=style)
                quarter_label.render(root)
                xy_cursor = (year_lenght_px * (index + 1) + year_lenght_px * 3 / 4, swimlane_height * 0.75)
                quarter_label = Rectangle(f"Q4",
                                          x=xy_cursor[0],
                                          y=xy_cursor[1],
                                          width=year_lenght_px / 4,
                                          height=swimlane_height * 0.25,
                                          style=style)
                quarter_label.render(root)



    def append_layers(self, root):
        # back to front order, lowest layer first
        layers = {'default': drawio_utils.create_layer(name='Default')}
        for layer in layers.values():
            root.append(layer)
        return root




class PowerPointRoadmapRenderer:
    def __init__(self):
        self.pptx_available = False
        self.Presentation = None
        self.Inches = None
        self.Pt = None
        self.RGBColor = None
        self.MSO_CONNECTOR = None

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.shapes import MSO_CONNECTOR

            self.pptx_available = True
            self.Presentation = Presentation
            self.Inches = Inches
            self.Pt = Pt
            self.RGBColor = RGBColor
            self.MSO_CONNECTOR = MSO_CONNECTOR
        except ImportError:
            pass

    def render(self, roadmap):
        if not self.pptx_available:
            print("python-pptx library is not installed. Please install it to use the PowerPoint renderer.")
            return "PowerPoint roadmap generation failed. Missing dependencies."

        prs = self.Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide layout

        # Add a text box for the title
        title_textbox = slide.shapes.add_textbox(self.Inches(0.5), self.Inches(0.5), self.Inches(9), self.Inches(1))
        title_textbox.text = "Roadmap"
        title_textbox.text_frame.paragraphs[0].font.size = self.Pt(24)
        title_textbox.text_frame.paragraphs[0].font.bold = True

        year_length_px = self.Inches(2)
        swimlane_height = self.Inches(0.5)
        start_x = self.Inches(0.5)
        start_y = self.Inches(1.5)

        years = [str(x) for x in range(roadmap.start_year, roadmap.start_year + roadmap.years)]

        for index, swimlane in enumerate(roadmap.swimlanes):
            slide.shapes.add_textbox(start_x, start_y + index * swimlane_height, year_length_px, swimlane_height)
            textbox = slide.shapes.add_textbox(start_x, start_y + index * swimlane_height, year_length_px,
                                               swimlane_height)
            textbox.text = swimlane.name
            textbox.text_frame.paragraphs[0].font.size = self.Pt(12)

            timeline_start_x = start_x + year_length_px
            timeline_end_x = timeline_start_x + year_length_px * len(years)
            timeline_y = start_y + index * swimlane_height + swimlane_height / 2

            line = slide.shapes.add_connector(self.MSO_CONNECTOR.STRAIGHT, timeline_start_x, timeline_y, timeline_end_x,
                                              timeline_y)
            line.line.color.rgb = self.RGBColor(0, 0, 0)

            for event in swimlane.events:
                event.render(year_length_px, roadmap.years)

        prs.save("roadmap.pptx")
        return "PowerPoint roadmap generated successfully!"
