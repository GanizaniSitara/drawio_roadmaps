from drawio_roadmaps.renderers.roadmap_renderer import RoadmapRenderer


class PowerPointRoadmapRenderer(RoadmapRenderer):
    def __init__(self):
        self.pptx_available = False
        self.Presentation = None
        self.Inches = None
        self.Pt = None
        self.RGBColor = None
        self.MSO_CONNECTOR = None

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.shapes import MSO_CONNECTOR

            self.pptx_available = True
            self.Presentation = Presentation
            self.Inches = Inches
            self.Pt = Pt
            self.RGBColor = RGBColor
            self.MSO_CONNECTOR = MSO_CONNECTOR
        except ImportError:
            pass

    def render(self, roadmap):
        if not self.pptx_available:
            print("python-pptx library is not installed. Please install it to use the PowerPoint renderer.")
            return "PowerPoint roadmap generation failed. Missing dependencies."

        prs = self.Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide layout

        # Add a text box for the title
        title_textbox = slide.shapes.add_textbox(self.Inches(0.5), self.Inches(0.5), self.Inches(9), self.Inches(1))
        title_textbox.text = "Roadmap"
        title_textbox.text_frame.paragraphs[0].font.size = self.Pt(24)
        title_textbox.text_frame.paragraphs[0].font.bold = True

        year_length_px = self.Inches(2)
        swimlane_height = self.Inches(0.5)
        start_x = self.Inches(0.5)
        start_y = self.Inches(1.5)

        years = [str(x) for x in range(roadmap.start_year, roadmap.start_year + roadmap.years)]

        for index, swimlane in enumerate(roadmap.swimlanes):
            slide.shapes.add_textbox(start_x, start_y + index * swimlane_height, year_length_px, swimlane_height)
            textbox = slide.shapes.add_textbox(start_x, start_y + index * swimlane_height, year_length_px,
                                               swimlane_height)
            textbox.text = swimlane.name
            textbox.text_frame.paragraphs[0].font.size = self.Pt(12)

            timeline_start_x = start_x + year_length_px
            timeline_end_x = timeline_start_x + year_length_px * len(years)
            timeline_y = start_y + index * swimlane_height + swimlane_height / 2

            line = slide.shapes.add_connector(self.MSO_CONNECTOR.STRAIGHT, timeline_start_x, timeline_y, timeline_end_x,
                                              timeline_y)
            line.line.color.rgb = self.RGBColor(0, 0, 0)

            for event in swimlane.events:
                event.render(year_length_px, roadmap.years)

        prs.save("roadmap.pptx")
        return "PowerPoint roadmap generated successfully!"
